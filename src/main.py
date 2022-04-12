import socket

from pptx import Presentation
import os
from flask import Flask, send_from_directory

from flask import request

gerarArquivo = Flask("name")


# pip install python-pptx


@gerarArquivo.route("/")
def gerar():
    try:
        return "<p>sucesso ao entrar</p>"
    except:
        return "<p>erro</p>"


def adicionarSlide(prs, titulo, estrofe):
    blank_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(blank_slide_layout)
    title = slide.shapes.title
    subtitulo = slide.placeholders[1]
    title.text = titulo
    subtitulo.text = estrofe.replace("<br>", '\n')


@gerarArquivo.route("/gerar", methods=['POST'])
def gerar_slides():
    # metodo para gerar e salvar o slides com a letra
    # recebendo os dados via POST
    global estrofe_slide
    textos = request.form.get('textos')
    titulo = request.form.get('titulo')
    modelo = request.form.get('modelo')
    try:
        # dividindo o texto em um lista apartir da seguinte separacao
        textoDividido = textos.split("</p>")
        # removendo o ultimo index da lista pois o mesmo fica vazio
        textoDividido.pop(-1)
        if 'geral' in modelo:
            caminho = "modelos_slides/modelo_geral.pptx"
            prs = Presentation(caminho)
        else:
            caminho = "modelos_slides/modelo_geracao_fire.pptx"
            prs = Presentation(caminho)
        for estrofe in textoDividido:  # pegando cada item da lista
            # divindo a lista na quantidade de vezes passada como parametro para o split
            dividir_estrofe = str(estrofe).split("<br>", 2)
            # verificando se a variavel atende quais situacoes
            if len(dividir_estrofe) == 1:
                estrofe_slide = dividir_estrofe[0]
            elif len(dividir_estrofe) >= 2:
                # definindo que a variavel vai receber o valor da variavel
                # pegando o primeiro index e o segundo index
                estrofe_slide = dividir_estrofe[0] + "\n" + dividir_estrofe[1]
            # adicionando um slide contendo os seguintes conteudos
            adicionarSlide(prs, titulo, estrofe_slide)
            # verificando se a variavel contem mais de dois index
            if len(dividir_estrofe) > 2:
                # divindo a lista na quantidade de vezes passada como
                # parametro para o split e atribuindo a uma variavel
                dividir_estrofe_secundaria = str(dividir_estrofe[2]).split("<br>", 2)
                # verificando se a variavel atende quais situacoes
                if len(dividir_estrofe_secundaria) == 1:
                    estrofe_slide_secundaria = dividir_estrofe_secundaria[0]
                    adicionarSlide(prs, titulo, estrofe_slide_secundaria)

                elif len(dividir_estrofe_secundaria) >= 2:
                    # definindo que a variavel vai receber o valor da variavel
                    # pegando o primeiro index e o segundo index
                    estrofe_slide_secundaria = dividir_estrofe_secundaria[0] + "\n" + dividir_estrofe_secundaria[1]
                    adicionarSlide(prs, titulo, estrofe_slide_secundaria)
                    # verificando se a variavel atende a seguinte situacao
                    if len(dividir_estrofe_secundaria) >= 3:
                        adicionarSlide(prs, titulo, dividir_estrofe_secundaria[2])

        prs.save(titulo + ".pptx")  # salvando o arquivo
        return "<p>sucesso</p>"
    except:
        return "<p>erro</p>"


def obterIP():
    # metodo para obter ip da maquina
    try:
        # obtendo ip
        ip = socket.gethostbyname(socket.gethostname())
        print(ip)
        # salvando o arquivo
        return ip
    except:
        return "<p>erro</p>"


@gerarArquivo.route("/baixarArquivo/<nome_arquivo>", methods=['GET'])
def baixar_arquivo(nome_arquivo):
    try:
        caminho_absoluto_arquivo_python = os.path.abspath(__file__)
        diretorio_src = os.path.dirname(caminho_absoluto_arquivo_python)
        print(diretorio_src)
        arquivo = nome_arquivo + ".pptx"
        return send_from_directory(diretorio_src,
                                   arquivo, as_attachment=True)
    except:
        return "<p>erro</p>"


@gerarArquivo.route("/excluirArquivo", methods=['POST'])
def excluir_arquivo():
    # metodo para excluir o arquivo gerado
    arquivo = request.form.get('arquivo')
    nome_arquivo = arquivo + ".pptx"
    caminho_absoluto_arquivo_python = os.path.abspath(__file__)
    diretorio_src = os.path.dirname(caminho_absoluto_arquivo_python)
    diretorio = os.listdir(diretorio_src)
    try:
        for file in diretorio:
            if file == nome_arquivo:
                os.remove(file)
        return "<p>sucesso ao excluir</p>"
    except:
        return "<p>erro ao excluir</p>"


if __name__ == '__main__':
    gerarArquivo.run(host=obterIP(), debug=True)
